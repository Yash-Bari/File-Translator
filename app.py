import streamlit as st
import os
from docx import Document
from pptx import Presentation
import PyPDF2
from deep_translator import GoogleTranslator
from base64 import b64encode

def translate_text(text, target_language):
    translated_text = GoogleTranslator(source='auto', target=target_language).translate(text)
    return translated_text

def translate_document(file, target_language):
    if file.type == "application/pdf":
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num] 
            text += page.extract_text()
        translated_text = translate_text(text, target_language)
        return translated_text
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = Document(file)
        text = "\n".join([p.text for p in doc.paragraphs])
        translated_text = translate_text(text, target_language)
        return translated_text
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        translated_text = translate_text(text, target_language)
        return translated_text
    else:
        st.error("Unsupported file format")

def main():
    st.title("Document Translator")

    uploaded_file = st.file_uploader("Upload a document", type=["docx", "pdf", "pptx"])

    if uploaded_file is not None:
        supported_languages = {
    'English': 'en',
    'French': 'fr',
    'Spanish': 'es',
    'German': 'de',
    'Chinese (Simplified)': 'zh-CN',
    'Japanese': 'ja',
    'Korean': 'ko',
    'Russian': 'ru',
    'Arabic': 'ar',
    'Italian': 'it',
    'Hindi': 'hi',
    'Marathi': 'mr',
    'Portuguese': 'pt',
    'Dutch': 'nl',
    'Swedish': 'sv',
    'Greek': 'el',
    'Thai': 'th',
    'Turkish': 'tr',
    'Vietnamese': 'vi',
    'Indonesian': 'id'
}


        target_language = st.selectbox("Select Target Language:", list(supported_languages.keys()))

        target_language_code = supported_languages[target_language]

        translated_text = translate_document(uploaded_file, target_language_code)

        st.subheader("Translated Text:")
        with st.expander("Translated Text", expanded=True):
            st.text_area(label="", value=translated_text, height=300)

        if st.button("Download Translated Document"):
            translated_filename = f"translated_{os.path.splitext(uploaded_file.name)[0]}_{target_language}.txt"
            with open(translated_filename, "w", encoding="utf-8") as f:
                f.write(translated_text)
            with open(translated_filename, "rb") as f:
                bin_data = f.read()
            data_url = f"data:application/octet-stream;base64,{b64encode(bin_data).decode()}"
            st.markdown(f'<a href="{data_url}" download="{translated_filename}"><button style="background-color: #4CAF50; color: white; padding: 10px 24px; border: none; cursor: pointer; border-radius: 5px;">Download {translated_filename}</button></a>', unsafe_allow_html=True)

if __name__ == "__main__":
    main()
